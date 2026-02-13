#!/usr/bin/env python3
"""
Script to create a comprehensive educational PowerPoint presentation
on Medical Image Segmentation using Attention U-Net
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Create the comprehensive medical image segmentation presentation"""
    
    # Create presentation with widescreen layout
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme - professional dark navy blue theme
    NAVY_BLUE = RGBColor(25, 55, 100)
    LIGHT_BLUE = RGBColor(70, 130, 180)
    WHITE = RGBColor(255, 255, 255)
    ACCENT_ORANGE = RGBColor(255, 140, 0)
    LIGHT_GRAY = RGBColor(220, 220, 220)
    
    def add_title_slide(title, subtitle=""):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = LIGHT_GRAY
        
        return slide
    
    def add_content_slide(title, content_items, notes="", slide_number=None):
        """Add a content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Divider line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_ORANGE
        line.line.width = Pt(3)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.3))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            
            # Check if item is a tuple (for nested bullets or different formatting)
            if isinstance(item, tuple):
                text, level = item
                p.text = text
                p.level = level
            else:
                p.text = item
                p.level = 0
            
            p.font.size = Pt(20) if p.level == 0 else Pt(18)
            p.font.color.rgb = WHITE
            p.space_before = Pt(10)
            p.space_after = Pt(6)
        
        # Slide number
        if slide_number:
            number_box = slide.shapes.add_textbox(Inches(9.2), Inches(7), Inches(0.5), Inches(0.3))
            number_frame = number_box.text_frame
            number_frame.text = str(slide_number)
            number_para = number_frame.paragraphs[0]
            number_para.font.size = Pt(14)
            number_para.font.color.rgb = LIGHT_GRAY
            number_para.alignment = PP_ALIGN.RIGHT
        
        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
        
        return slide
    
    slide_num = 0
    
    # ==================== TITLE SLIDE ====================
    add_title_slide(
        "Medical Image Segmentation using Deep Learning",
        "From U-Net to Attention U-Net: A Comprehensive Overview"
    )
    
    # ==================== SECTION 1: INTRODUCTION ====================
    slide_num += 1
    add_content_slide(
        "What is Medical Image Segmentation?",
        [
            "• Definition: The process of partitioning medical images into meaningful regions",
            "• Goal: Precisely delineate anatomical structures and pathological regions",
            "• Output: Pixel-wise classification of tissue types, organs, or lesions",
            "",
            "• Critical for:",
            ("  - Diagnosis and treatment planning", 1),
            ("  - Surgical navigation and intervention", 1),
            ("  - Quantitative analysis (volume, shape, location)", 1),
            ("  - Monitoring disease progression and treatment response", 1)
        ],
        notes="Medical image segmentation is a fundamental task in medical image analysis. Unlike classification which assigns a single label to an entire image, segmentation provides spatial information by classifying every pixel. This is essential for understanding the precise location, size, and shape of anatomical structures.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Clinical Applications",
        [
            "• Lung Cancer Screening & Diagnosis",
            ("  - Lung nodule detection and segmentation", 1),
            ("  - Tumor volume quantification for staging", 1),
            ("  - Treatment response assessment", 1),
            "",
            "• Brain Tumor Analysis",
            ("  - Glioma segmentation (enhancing, non-enhancing, edema)", 1),
            ("  - Surgical planning and resection guidance", 1),
            ("  - Radiotherapy target delineation", 1),
            "",
            "• Organ Segmentation",
            ("  - Liver, kidney, heart segmentation for surgical planning", 1),
            ("  - Cardiac chamber quantification for functional assessment", 1),
            ("  - Multi-organ segmentation for radiation therapy planning", 1)
        ],
        notes="Medical image segmentation has wide-ranging clinical applications. In lung cancer, accurate segmentation helps radiologists measure tumor size and track changes over time. For brain tumors, different tumor components need to be segmented separately. Organ segmentation is crucial for transplant planning, surgical guidance, and ensuring radiation therapy targets tumors while sparing healthy tissue.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "The Challenge: Why Manual Segmentation Fails",
        [
            "• Manual segmentation is the current gold standard BUT:",
            "",
            "• Time-Consuming",
            ("  - Single 3D MRI scan: 1-4 hours per expert radiologist", 1),
            ("  - Hospital burden: thousands of scans per year", 1),
            "",
            "• Inter-Observer Variability",
            ("  - Different radiologists produce different segmentations", 1),
            ("  - Inconsistent boundaries and measurements", 1),
            "",
            "• Subjective & Error-Prone",
            ("  - Human fatigue affects accuracy", 1),
            ("  - Subtle pathologies may be missed", 1),
            "",
            "• Not Scalable",
            ("  - Cannot meet increasing demand for medical imaging", 1),
            ("  - Limited availability of expert radiologists", 1)
        ],
        notes="Manual segmentation is extremely labor-intensive and requires expert knowledge. A single 3D brain MRI can take 2-4 hours to segment manually. Inter-observer variability is a major issue - studies show Dice scores between different experts ranging from 0.7-0.9, indicating significant disagreement. This motivates the need for automated, consistent, and fast segmentation algorithms.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Why Deep Learning?",
        [
            "• Advantages over Traditional Methods:",
            "",
            "• End-to-End Learning",
            ("  - Learns features automatically from data", 1),
            ("  - No hand-crafted feature engineering required", 1),
            "",
            "• Superior Performance",
            ("  - Matches or exceeds human expert performance", 1),
            ("  - Handles complex, high-dimensional data", 1),
            "",
            "• Speed & Consistency",
            ("  - Inference in seconds instead of hours", 1),
            ("  - Reproducible results with zero inter-observer variability", 1),
            "",
            "• Continuous Improvement",
            ("  - Performance improves with more training data", 1),
            ("  - Transfer learning leverages pre-trained models", 1)
        ],
        notes="Deep learning has revolutionized medical image segmentation. Unlike traditional methods that require manual feature engineering (e.g., designing filters to detect edges), deep learning models learn optimal features directly from data. Modern architectures like U-Net achieve Dice scores >0.9 on many tasks, rivaling expert radiologists while processing images in seconds.",
        slide_number=slide_num
    )
    
    # ==================== SECTION 2: EVOLUTION OF AI ====================
    slide_num += 1
    add_content_slide(
        "Traditional Methods (Pre-2015)",
        [
            "• Thresholding & Region Growing",
            ("  - Simple pixel intensity-based methods", 1),
            ("  - Limited to homogeneous regions, sensitive to noise", 1),
            "",
            "• Active Contours (Snakes)",
            ("  - Iteratively deform contour to match boundaries", 1),
            ("  - Requires good initialization, local minima issues", 1),
            "",
            "• Atlas-Based Segmentation",
            ("  - Register pre-segmented atlas to new image", 1),
            ("  - Limited by registration accuracy and anatomical variation", 1),
            "",
            "• Limitations:",
            ("  - Required domain expertise for feature engineering", 1),
            ("  - Poor generalization to different imaging modalities", 1),
            ("  - Failed on complex, heterogeneous pathologies", 1)
        ],
        notes="Before deep learning, segmentation relied on classical computer vision techniques. Thresholding separates pixels based on intensity values but fails when regions overlap in intensity. Active contours deform a curve to match object boundaries but get stuck in local optima. Atlas-based methods work well for healthy anatomy but struggle with pathology that deforms normal structure.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "FCN - Fully Convolutional Networks (2015)",
        [
            "• Long, J., Shelhamer, E., & Darrell, T. (2015)",
            "  'Fully Convolutional Networks for Semantic Segmentation'",
            "  CVPR 2015",
            "",
            "• Key Innovation: Replace fully connected layers with convolutions",
            ("  - Accepts arbitrary input sizes", 1),
            ("  - Outputs spatial segmentation maps", 1),
            ("  - Uses skip connections to combine coarse and fine features", 1),
            "",
            "• Architecture: Encoder (VGG) + Decoder (transposed convolutions)",
            "",
            "• Results: PASCAL VOC 2012 - 62.2% mean IU",
            "",
            "• Impact: Pioneered end-to-end learning for dense prediction",
        ],
        notes="FCN by Long et al. was the first work to adapt CNNs for semantic segmentation. The key insight was replacing fully connected layers with 1x1 convolutions, allowing the network to output spatial maps. They introduced skip connections to combine high-level semantic information with low-level spatial details. This laid the foundation for all modern segmentation architectures.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "U-Net - The Breakthrough (2015)",
        [
            "• Ronneberger, O., Fischer, P., & Brox, T. (2015)",
            "  'U-Net: Convolutional Networks for Biomedical Image Segmentation'",
            "  MICCAI 2015",
            "",
            "• Key Innovation: Symmetric encoder-decoder with skip connections",
            ("  - Contracting path (encoder): captures context", 1),
            ("  - Expanding path (decoder): enables precise localization", 1),
            ("  - Skip connections: preserve spatial information", 1),
            "",
            "• Design for Medical Imaging:",
            ("  - Works with very few training images (30-50)", 1),
            ("  - Data augmentation strategy for small datasets", 1),
            "",
            "• Results: ISBI Cell Segmentation - Won by large margin",
            "  (IOU = 92%, far exceeding previous methods)",
        ],
        notes="U-Net revolutionized medical image segmentation. The symmetric U-shaped architecture with skip connections became the de facto standard. Skip connections concatenate encoder features with decoder features at the same resolution, enabling the network to use both high-resolution spatial information and high-level semantic information. Critically, U-Net was designed to work with small datasets common in medical imaging.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "V-Net - 3D Extension (2016)",
        [
            "• Milletari, F., Navab, N., & Ahmadi, S. A. (2016)",
            "  'V-Net: Fully Convolutional Neural Networks for Volumetric Medical Image Segmentation'",
            "  3DV 2016",
            "",
            "• Key Innovations:",
            ("  - Full 3D convolutions for volumetric segmentation", 1),
            ("  - Introduced Dice loss function for training", 1),
            ("  - Residual connections within conv blocks", 1),
            "",
            "• Architecture: 3D U-Net style with residual blocks",
            "",
            "• Dice Loss: Directly optimizes segmentation overlap metric",
            "  Dice = 2|X ∩ Y| / (|X| + |Y|)",
            "",
            "• Results: Prostate MRI segmentation - Dice 0.869",
        ],
        notes="V-Net extended U-Net to 3D volumes, essential for medical imaging where data is inherently volumetric (CT, MRI). The Dice loss was a major contribution - instead of pixel-wise cross-entropy, it directly optimizes the Dice coefficient metric that radiologists use. This helps with class imbalance (most pixels are background) and directly maximizes the evaluation metric.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "U-Net++ - Nested Architecture (2018)",
        [
            "• Zhou, Z., Rahman Siddiquee, M. M., Tajbakhsh, N., & Liang, J. (2018)",
            "  'UNet++: A Nested U-Net Architecture for Medical Image Segmentation'",
            "  MICCAI Workshop 2018",
            "",
            "• Key Innovation: Nested skip pathways with dense connections",
            ("  - Multiple U-Nets of different depths nested together", 1),
            ("  - Dense skip connections reduce semantic gap", 1),
            ("  - Deep supervision at multiple scales", 1),
            "",
            "• Architecture: X1,0, X2,0, ..., X5,0 in encoder",
            "              X0,1, X0,2, ..., X0,5 nested nodes with dense connections",
            "",
            "• Results:",
            ("  - Liver: Dice 0.826 (U-Net: 0.815)", 1),
            ("  - Lung nodules: Dice 0.838 (U-Net: 0.820)", 1),
        ],
        notes="U-Net++ addresses the semantic gap between encoder and decoder features in U-Net. By introducing nested dense skip connections, features are gradually upsampled and refined. Deep supervision means auxiliary losses are applied at intermediate layers, helping gradient flow. This leads to better performance but at the cost of more parameters and computation.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Attention U-Net (2018)",
        [
            "• Oktay, O., et al. (2018)",
            "  'Attention U-Net: Learning Where to Look for the Pancreas'",
            "  MIDL 2018",
            "",
            "• Key Innovation: Attention gates in skip connections",
            ("  - Learn to focus on salient features", 1),
            ("  - Suppress irrelevant background regions", 1),
            ("  - No additional supervision required", 1),
            "",
            "• Attention Mechanism:",
            "  α = σ(ψ(ReLU(W_g * g + W_x * x)))",
            "  where g = gating signal (decoder), x = skip connection (encoder)",
            "",
            "• Results:",
            ("  - Pancreas CT: Dice 0.859 (U-Net: 0.843)", 1),
            ("  - Liver tumor: Dice 0.681 (U-Net: 0.658)", 1),
        ],
        notes="Attention U-Net introduces attention gates that act as a soft selection mechanism. The gating signal from the decoder pathway guides the attention mechanism to focus on specific regions in the encoder features. This is particularly useful in medical imaging where the region of interest is small compared to the entire image (e.g., small tumors, pancreas). The attention weights are learned end-to-end without extra supervision.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "nnU-Net - Self-Configuring Framework (2021)",
        [
            "• Isensee, F., Jaeger, P. F., Kohl, S. A., Petersen, J., & Maier-Hein, K. H. (2021)",
            "  'nnU-Net: a self-configuring method for deep learning-based biomedical image segmentation'",
            "  Nature Methods, 18(2), 203-211",
            "",
            "• Key Innovation: Automatic configuration based on dataset properties",
            ("  - Automatically determines preprocessing, architecture, training", 1),
            ("  - No manual tuning required", 1),
            ("  - Robust across diverse medical imaging tasks", 1),
            "",
            "• Three U-Net variants: 2D, 3D full-res, 3D cascade",
            "",
            "• Results: Won 33 out of 53 international challenges",
            ("  - Medical Segmentation Decathlon: 1st place", 1),
            ("  - KiTS, BraTS, LiTS challenges: Top performance", 1),
        ],
        notes="nnU-Net by Isensee et al. represents the engineering pinnacle of U-Net based methods. It automatically configures all preprocessing, architecture choices, and training hyperparameters based on the dataset fingerprint. This removes the need for manual tuning and makes it highly robust. It has become the de facto baseline for medical image segmentation competitions and research.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Transformer-Based Approaches (2021-2023)",
        [
            "• TransUNet (Chen et al., 2021)",
            ("  - Combines CNN encoder with Transformer decoder", 1),
            ("  - Self-attention captures global context", 1),
            ("  - Synapse multi-organ CT: Dice 0.772 (U-Net: 0.688)", 1),
            "",
            "• Swin-UNET (Cao et al., 2022)",
            ("  - Pure transformer architecture with Swin Transformer blocks", 1),
            ("  - Hierarchical shifted windows for efficiency", 1),
            ("  - Competitive with CNN-based methods", 1),
            "",
            "• nnFormer (Zhou et al., 2023)",
            ("  - Interleaved local attention and global attention", 1),
            ("  - BTCV: Dice 0.873 (best transformer method)", 1),
            "",
            "• Limitation: Requires large datasets (>10k images) to outperform CNNs",
        ],
        notes="Transformers, which revolutionized NLP, have been adapted for medical image segmentation. The self-attention mechanism can capture long-range dependencies that CNNs miss. However, transformers require much more data to train effectively. On large datasets, they can outperform CNNs, but on small medical imaging datasets (hundreds of images), CNN-based methods like nnU-Net still dominate.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Evolution Timeline",
        [
            "2015: FCN - First end-to-end segmentation network",
            "      U-Net - Symmetric encoder-decoder for medical imaging ⭐",
            "",
            "2016: V-Net - 3D extension + Dice loss",
            "",
            "2018: U-Net++ - Nested dense skip connections",
            "      Attention U-Net - Attention gates ⭐",
            "",
            "2021: nnU-Net - Self-configuring framework ⭐",
            "      TransUNet - Transformers meet U-Net",
            "",
            "2022: Swin-UNET - Pure transformer architecture",
            "",
            "2023: nnFormer - Efficient transformer",
            "      Foundation Models - SAM, MedSAM (emerging)",
        ],
        notes="This timeline shows the rapid evolution of medical image segmentation. U-Net (2015) remains the foundation. Attention U-Net (2018) added interpretability and performance gains. nnU-Net (2021) proved that proper engineering of U-Net variants can beat more complex methods. Transformers are emerging but haven't fully replaced CNNs yet. Foundation models like SAM show promise for universal segmentation.",
        slide_number=slide_num
    )
    
    # ==================== SECTION 3: ATTENTION U-NET DEEP DIVE ====================
    slide_num += 1
    add_content_slide(
        "Attention U-Net: Motivation",
        [
            "• Problem with Standard U-Net Skip Connections:",
            ("  - All encoder features are concatenated equally", 1),
            ("  - No mechanism to emphasize relevant regions", 1),
            ("  - Background features can overwhelm small target regions", 1),
            "",
            "• Attention U-Net Solution:",
            ("  - Add attention gates before concatenation", 1),
            ("  - Learn to highlight salient features", 1),
            ("  - Suppress irrelevant activations", 1),
            "",
            "• Benefits:",
            ("  - Improved sensitivity for small structures", 1),
            ("  - Better handling of class imbalance", 1),
            ("  - Interpretability through attention weights", 1),
            ("  - Minimal computational overhead", 1),
        ],
        notes="Standard U-Net skip connections concatenate all encoder features without discrimination. This can be problematic when the target region is small (e.g., small tumors, pancreas) compared to background. Attention gates learn to weight features based on their relevance, essentially teaching the network where to look. This improves both performance and interpretability.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Attention Gate Mechanism",
        [
            "• Inputs:",
            ("  - x: Skip connection features from encoder (high resolution)", 1),
            ("  - g: Gating signal from decoder (coarse, semantic)", 1),
            "",
            "• Operation:",
            "  1. Transform x and g to same channel dimension:",
            "     W_x(x) + W_g(g)  [element-wise addition after 1×1 conv]",
            "",
            "  2. Apply ReLU activation",
            "",
            "  3. Generate attention coefficients:",
            "     α = σ(ψ(ReLU(W_x(x) + W_g(g))))",
            "     where ψ is 1×1 conv, σ is sigmoid",
            "",
            "  4. Apply attention to skip connection:",
            "     x̂ = α ⊙ x  [element-wise multiplication]",
        ],
        notes="The attention gate takes two inputs: high-resolution features from encoder and coarse semantic features from decoder. These are transformed and combined, then passed through a sigmoid to produce attention coefficients between 0 and 1. These coefficients weight each spatial location and channel in the encoder features, suppressing background and highlighting foreground.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Attention Gate: Mathematical Formulation",
        [
            "• Full Mathematical Expression:",
            "",
            "  q_att = ψ^T(σ_1(W_x^T x + W_g^T g + b_g)) + b_ψ",
            "",
            "  α = σ_2(q_att)",
            "",
            "  x̂ = x ⊙ α",
            "",
            "• Where:",
            ("  - W_x, W_g: Linear transformations (1×1 conv)", 1),
            ("  - ψ: Attention coefficient generator (1×1 conv)", 1),
            ("  - σ_1: ReLU activation", 1),
            ("  - σ_2: Sigmoid activation", 1),
            ("  - ⊙: Element-wise multiplication (broadcast)", 1),
            "",
            "• Key Property: Differentiable, trained end-to-end with backprop",
        ],
        notes="This mathematical formulation shows that attention gates are fully differentiable and learned through standard backpropagation. The transformations W_x and W_g are 1x1 convolutions that project features to an intermediate dimension. The gating signal g provides context about what the network is trying to segment, which guides the attention. No additional labels or supervision is needed - the attention learns where to focus based on the segmentation task.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Attention U-Net vs Standard U-Net",
        [
            "• Architecture Differences:",
            ("  - Standard U-Net: Direct concatenation of skip connections", 1),
            ("  - Attention U-Net: Attention gate → weighted features → concat", 1),
            "",
            "• Performance Comparison (Pancreas CT Segmentation):",
            ("  - U-Net: Dice = 0.843, Sensitivity = 0.851", 1),
            ("  - Attention U-Net: Dice = 0.859, Sensitivity = 0.889", 1),
            ("  - Improvement: +1.6% Dice, +3.8% Sensitivity", 1),
            "",
            "• Computational Cost:",
            ("  - Parameters: ~2-3% increase", 1),
            ("  - Inference time: ~5% increase", 1),
            ("  - Negligible overhead for significant performance gain", 1),
            "",
            "• Interpretability: Attention maps visualize focus regions",
        ],
        notes="Attention U-Net provides consistent improvements over standard U-Net, especially for small target regions where sensitivity is critical. The attention maps also provide interpretability - we can visualize what the network focuses on. The overhead is minimal: a small increase in parameters (attention modules) and negligible increase in inference time. This makes it an efficient upgrade to standard U-Net.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Attention Visualization & Interpretability",
        [
            "• Attention Maps Show Network Focus:",
            ("  - Bright regions: high attention (foreground)", 1),
            ("  - Dark regions: low attention (background)", 1),
            "",
            "• Progressive Refinement Through Decoder:",
            ("  - Early layers: coarse attention on general region", 1),
            ("  - Later layers: fine-grained attention on boundaries", 1),
            "",
            "• Clinical Interpretability:",
            ("  - Helps clinicians understand AI decision-making", 1),
            ("  - Can identify when model focuses on wrong regions", 1),
            ("  - Builds trust in automated segmentation systems", 1),
            "",
            "• Failure Mode Detection:",
            ("  - Low attention on known pathology → potential miss", 1),
            ("  - High attention on normal tissue → false positive", 1),
        ],
        notes="Attention maps provide valuable interpretability, especially important in medical applications. By visualizing attention weights at different decoder stages, we can see how the network progressively refines its focus from coarse regions to precise boundaries. This interpretability helps clinicians trust and validate AI predictions, and can help identify failure modes before they cause clinical errors.",
        slide_number=slide_num
    )
    
    # ==================== SECTION 4: OUR IMPLEMENTATION ====================
    slide_num += 1
    add_content_slide(
        "Our Implementation: Dataset",
        [
            "• Medical Segmentation Decathlon - Task 06: Lung",
            ("  - Source: http://medicaldecathlon.com/", 1),
            ("  - Modality: CT scans", 1),
            ("  - Target: Lung cancer segmentation", 1),
            ("  - Format: NIfTI (.nii.gz) 3D volumes", 1),
            "",
            "• Dataset Statistics:",
            ("  - Total volumes: 63 training cases", 1),
            ("  - Image dimensions: 512×512×variable depth", 1),
            ("  - Classes: 2 (background + lung tumor)", 1),
            "",
            "• Data Split:",
            ("  - Training: 80% (~50 cases)", 1),
            ("  - Validation: 10% (~6 cases)", 1),
            ("  - Testing: 10% (~6 cases)", 1),
        ],
        notes="We use the Medical Segmentation Decathlon Lung dataset, a standardized benchmark for medical image segmentation. The dataset contains CT scans with lung cancer annotations. We work with the middle slice of each 3D volume for 2D segmentation. The 80/10/10 split ensures we have enough data for training while preserving independent validation and test sets.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Preprocessing Pipeline",
        [
            "• NIfTI Volume Loading (data/dataset.py):",
            ("  - Load 3D volumes using nibabel library", 1),
            ("  - Extract middle slice for 2D segmentation", 1),
            "",
            "• Intensity Normalization:",
            ("  - Min-max normalization to [0, 1] range", 1),
            ("  - Per-image normalization: (x - min) / (max - min + ε)", 1),
            "",
            "• Data Augmentation (Training Only):",
            ("  - Horizontal flip (p=0.5)", 1),
            ("  - Vertical flip (p=0.5)", 1),
            ("  - Rotation (±15°, p=0.5)", 1),
            ("  - Elastic deformation (p=0.5)", 1),
            ("  - Random brightness/contrast (p=0.5)", 1),
            ("  - Gaussian noise (p=0.5)", 1),
            "",
            "• Tensor Conversion: NumPy → PyTorch tensors (C, H, W)",
        ],
        notes="Our preprocessing pipeline is implemented in data/dataset.py. We normalize intensities per-image since CT scans can have varying intensity ranges. Augmentation is critical for medical imaging where training data is limited. We use Albumentations library for efficient augmentations. Elastic deformation is particularly important as it simulates realistic anatomical variations.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Model Architecture Details",
        [
            "• Attention U-Net (model/unet.py):",
            ("  - Input: 1 channel (grayscale CT)", 1),
            ("  - Output: 2 classes (background + tumor)", 1),
            ("  - Base channels: 32", 1),
            "",
            "• Encoder (Contracting Path):",
            ("  - 5 levels: 32 → 64 → 128 → 256 → 512 channels", 1),
            ("  - Each level: 2 conv blocks with residual connections", 1),
            ("  - MaxPool 2×2 for downsampling", 1),
            ("  - Dropout (p=0.1) for regularization", 1),
            "",
            "• Decoder (Expanding Path):",
            ("  - 4 levels: 512 → 256 → 128 → 64 → 32 channels", 1),
            ("  - Bilinear upsampling (2×)", 1),
            ("  - Attention gates at each level", 1),
            ("  - Concatenation with attended encoder features", 1),
            "",
            "• Output: 1×1 conv to produce 2-channel prediction",
        ],
        notes="Our implementation in model/unet.py follows the Attention U-Net architecture with some enhancements. We use residual connections within conv blocks for better gradient flow. The base channels of 32 gives us ~4M parameters - a good balance for our dataset size. Dropout provides regularization to prevent overfitting. Attention gates are applied at all 4 skip connections.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Deep Supervision Strategy",
        [
            "• Deep Supervision: Auxiliary losses at intermediate layers",
            "",
            "• Implementation (model/unet.py:151-181):",
            ("  - Main output: final decoder output (full resolution)", 1),
            ("  - Auxiliary outputs: 3 intermediate decoder layers", 1),
            ("  - Each has 1×1 conv to produce class predictions", 1),
            "",
            "• Training Loss:",
            "  L_total = L_main + 0.25 × (L_aux1 + L_aux2 + L_aux3)",
            "",
            "• Benefits:",
            ("  - Improved gradient flow to deep layers", 1),
            ("  - Regularization effect (prevents overfitting)", 1),
            ("  - Faster convergence", 1),
            "",
            "• Inference: Only main output used (deep_supervision=False)",
        ],
        notes="Deep supervision is a technique where we add auxiliary losses at intermediate decoder layers. During training, targets are downsampled to match auxiliary output resolutions. This provides direct supervision to intermediate layers, improving gradient flow. The auxiliary losses are weighted 0.25 relative to the main loss. During inference, we disable deep supervision and use only the final output for efficiency.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Training Configuration",
        [
            "• Loss Function (train.py:15-52):",
            ("  - Combined Dice + Cross-Entropy Loss", 1),
            ("  - Dice loss: 1 - (2|X∩Y| + ε) / (|X| + |Y| + ε)", 1),
            ("  - CE loss: Standard cross-entropy", 1),
            ("  - Weights: 0.5 Dice + 0.5 CE", 1),
            "",
            "• Optimizer: Adam (lr=1e-3)",
            "",
            "• Learning Rate Scheduler: CosineAnnealingLR",
            ("  - Smooth decay from 1e-3 to 1e-6", 1),
            ("  - T_max = num_epochs", 1),
            "",
            "• Mixed Precision Training (AMP):",
            ("  - FP16 for forward/backward, FP32 for updates", 1),
            ("  - 2× speedup with minimal accuracy loss", 1),
        ],
        notes="Our training strategy is implemented in train.py. The combined loss leverages both Dice (good for overlap) and CE (good for pixel-wise accuracy). Adam optimizer is standard for medical imaging. CosineAnnealingLR provides smooth learning rate decay. Mixed precision training using PyTorch AMP reduces memory usage and speeds up training by ~2× with negligible impact on final accuracy.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Training Details & Hyperparameters",
        [
            "• Hyperparameters (config.py):",
            ("  - Batch size: 4", 1),
            ("  - Epochs: 100", 1),
            ("  - Early stopping patience: 15 epochs", 1),
            ("  - Num workers: 4 (data loading)", 1),
            "",
            "• Regularization:",
            ("  - Dropout: 0.1", 1),
            ("  - Data augmentation (see preprocessing slide)", 1),
            ("  - Deep supervision", 1),
            "",
            "• Monitoring:",
            ("  - TensorBoard logging (loss, Dice, learning rate)", 1),
            ("  - Checkpoint saving (best model + periodic)", 1),
            ("  - Early stopping on validation Dice", 1),
        ],
        notes="We use a batch size of 4 due to memory constraints - medical images are high resolution. Training for 100 epochs with early stopping prevents overfitting while ensuring convergence. TensorBoard integration provides real-time monitoring of training progress. We save the best model based on validation Dice score and also periodic checkpoints every 10 epochs for analysis.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Evaluation Metrics",
        [
            "• Dice Similarity Coefficient (DSC):",
            "  Dice = 2|X ∩ Y| / (|X| + |Y|)",
            ("  - Measures overlap between prediction and ground truth", 1),
            ("  - Range: [0, 1], higher is better", 1),
            ("  - Standard metric in medical image segmentation", 1),
            "",
            "• Intersection over Union (IoU / Jaccard):",
            "  IoU = |X ∩ Y| / |X ∪ Y|",
            ("  - Related to Dice: Dice = 2×IoU / (1 + IoU)", 1),
            "",
            "• Hausdorff Distance (HD95):",
            ("  - Maximum distance from prediction to ground truth boundary", 1),
            ("  - Measures worst-case boundary error", 1),
            ("  - Lower is better (measured in pixels)", 1),
            "",
            "• Implementation: evaluate.py",
        ],
        notes="Dice coefficient is the most common metric in medical segmentation - it measures overlap and is robust to class imbalance. IoU is similar but more stringent. Hausdorff distance measures boundary accuracy, critical for surgical planning where precise boundaries matter. We implement all three metrics in evaluate.py for comprehensive assessment. Dice >0.8 is generally considered good, >0.9 is excellent.",
        slide_number=slide_num
    )
    
    # ==================== SECTION 5: RESULTS (PLACEHOLDERS) ====================
    slide_num += 1
    add_content_slide(
        "[PLACEHOLDER] Training Curves",
        [
            "[INSERT RESULTS AFTER TRAINING]",
            "",
            "• Training & Validation Loss over Epochs",
            ("  - Plot showing decreasing loss curves", 1),
            ("  - Convergence behavior", 1),
            ("  - Overfitting detection", 1),
            "",
            "• Dice Score over Epochs",
            ("  - Training and validation Dice progression", 1),
            ("  - Best validation Dice achieved", 1),
            ("  - Epoch where early stopping triggered", 1),
            "",
            "• Learning Rate Schedule",
            ("  - CosineAnnealing decay visualization", 1),
        ],
        notes="PLACEHOLDER: After training, insert line plots showing: (1) Training and validation loss curves over epochs, (2) Training and validation Dice score curves, (3) Learning rate schedule. These curves help diagnose convergence, overfitting, and training stability. Typical training takes 30-50 epochs to converge with early stopping.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "[PLACEHOLDER] Qualitative Results",
        [
            "[INSERT RESULTS AFTER TRAINING]",
            "",
            "• Visualization Grid: Input | Ground Truth | Prediction",
            "",
            "• Best Cases:",
            ("  - Examples where model achieves high Dice (>0.9)", 1),
            ("  - Clean tumor boundaries", 1),
            ("  - Correct size and shape", 1),
            "",
            "• Average Cases:",
            ("  - Typical performance (Dice 0.7-0.9)", 1),
            ("  - Minor boundary errors", 1),
            "",
            "• Challenging Cases:",
            ("  - Small tumors (high difficulty)", 1),
            ("  - Low contrast regions", 1),
            ("  - Noisy images", 1),
        ],
        notes="PLACEHOLDER: After training, create a grid visualization showing input CT slices, ground truth masks, and model predictions side-by-side. Include diverse examples: best cases (perfect segmentation), average cases (typical performance), and challenging cases (small tumors, low contrast, noise). This provides qualitative assessment of model performance and failure modes. Use visualizations/evaluation/ directory.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "[PLACEHOLDER] Quantitative Results",
        [
            "[INSERT RESULTS AFTER TRAINING]",
            "",
            "• Test Set Performance:",
            "",
            "  Metric                  | Mean ± Std",
            "  ----------------------- | -----------",
            "  Dice Score              | [TBD] ± [TBD]",
            "  IoU Score               | [TBD] ± [TBD]",
            "  Hausdorff Distance (mm) | [TBD] ± [TBD]",
            "",
            "• Per-Class Breakdown:",
            ("  - Background class: [TBD]", 1),
            ("  - Tumor class: [TBD]", 1),
            "",
            "• Inference Time: [TBD] ms per image",
        ],
        notes="PLACEHOLDER: After evaluation on test set, fill in quantitative metrics. Expected performance for lung tumor segmentation: Dice 0.7-0.85 (challenging task due to tumor heterogeneity), IoU 0.6-0.75, Hausdorff distance <10mm. Report mean and standard deviation across test set. Also report per-class metrics and inference time on your hardware (typically 20-50ms per image on GPU).",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "[PLACEHOLDER] Failure Cases & Analysis",
        [
            "[INSERT RESULTS AFTER TRAINING]",
            "",
            "• Common Failure Modes:",
            "",
            "1. Small Tumors",
            ("  - Dice: [TBD] (< mean)", 1),
            ("  - Cause: Limited training examples, class imbalance", 1),
            "",
            "2. Low Contrast Regions",
            ("  - Tumor intensity similar to surrounding tissue", 1),
            ("  - Model struggles with boundary delineation", 1),
            "",
            "3. Irregular Shapes",
            ("  - Non-spherical, heterogeneous tumors", 1),
            ("  - Over-smoothing or under-segmentation", 1),
            "",
            "• Mitigation Strategies:",
            ("  - More aggressive augmentation", 1),
            ("  - Focal loss for hard examples", 1),
            ("  - Ensemble methods", 1),
        ],
        notes="PLACEHOLDER: After evaluation, analyze failure cases to understand model limitations. Common failures in lung tumor segmentation: (1) Small tumors (<10mm) are often missed or under-segmented, (2) Low contrast tumors blend with surrounding tissue, (3) Irregular shapes are over-smoothed. Document specific examples with Dice scores. Suggest improvements: focal loss weights hard examples, test-time augmentation, multi-model ensembles.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "[PLACEHOLDER] Comparison with Baseline U-Net",
        [
            "[INSERT RESULTS AFTER TRAINING BOTH MODELS]",
            "",
            "• Performance Comparison:",
            "",
            "  Metric                  | U-Net | Attention U-Net | Δ",
            "  ----------------------- | ----- | --------------- | ---",
            "  Dice Score              | [TBD] | [TBD]           | [TBD]",
            "  IoU Score               | [TBD] | [TBD]           | [TBD]",
            "  Hausdorff Distance (mm) | [TBD] | [TBD]           | [TBD]",
            "  Parameters (M)          | [TBD] | [TBD]           | [TBD]",
            "  Inference Time (ms)     | [TBD] | [TBD]           | [TBD]",
            "",
            "• Expected Improvement: +1-3% Dice, +2-5% Sensitivity",
            "",
            "• Analysis: Is the attention mechanism beneficial?",
        ],
        notes="PLACEHOLDER: To validate the attention mechanism, train both standard U-Net and Attention U-Net on the same data. Compare Dice, IoU, and Hausdorff distance. Also compare parameter count and inference time. Expected results: Attention U-Net should show 1-3% Dice improvement with minimal computational overhead. Attention is especially beneficial for small tumors where it improves sensitivity.",
        slide_number=slide_num
    )
    
    # ==================== SECTION 6: CONCLUSION ====================
    slide_num += 1
    add_content_slide(
        "Summary of Findings",
        [
            "• Medical Image Segmentation is Critical",
            ("  - Essential for diagnosis, treatment planning, and monitoring", 1),
            ("  - Manual segmentation is slow, inconsistent, and non-scalable", 1),
            "",
            "• Deep Learning Has Revolutionized the Field",
            ("  - U-Net (2015) established the foundation", 1),
            ("  - Attention mechanisms improve focus on salient regions", 1),
            ("  - Modern methods (nnU-Net) match human expert performance", 1),
            "",
            "• Our Implementation: Attention U-Net for Lung Tumors",
            ("  - Dataset: Medical Segmentation Decathlon - Lung CT", 1),
            ("  - Combined Dice + CE loss, deep supervision, AMP training", 1),
            ("  - Comprehensive evaluation: Dice, IoU, Hausdorff distance", 1),
            "",
            "• Key Takeaway: Attention gates provide interpretability + performance",
        ],
        notes="This lecture covered the evolution of medical image segmentation from traditional methods to modern deep learning approaches. U-Net revolutionized the field, and Attention U-Net improved it with learnable focus mechanisms. Our implementation demonstrates a complete pipeline from data loading to evaluation. The key innovation of attention gates is providing both performance gains and interpretability, critical for clinical deployment.",
        slide_number=slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "Future Directions & Open Challenges",
        [
            "• 3D Attention Mechanisms",
            ("  - Current: 2D slices, Future: Full 3D volumetric attention", 1),
            ("  - Challenge: Memory and computational requirements", 1),
            "",
            "• Transformer-Based Architectures",
            ("  - Vision Transformers, Swin Transformers for segmentation", 1),
            ("  - Challenge: Requires large datasets (>10k images)", 1),
            "",
            "• Foundation Models",
            ("  - SAM (Segment Anything Model), MedSAM for medical imaging", 1),
            ("  - Transfer learning from massive datasets", 1),
            ("  - Challenge: Domain adaptation to medical imaging", 1),
            "",
            "• Multi-Modal Fusion",
            ("  - Combining CT, MRI, PET for comprehensive analysis", 1),
            "",
            "• Uncertainty Quantification",
            ("  - Bayesian deep learning, ensemble methods", 1),
            ("  - Critical for clinical decision support", 1),
        ],
        notes="The field continues to evolve rapidly. Future directions include: (1) 3D attention for volumetric data, (2) Transformer architectures once datasets grow larger, (3) Foundation models like SAM that can segment anything with minimal fine-tuning, (4) Multi-modal fusion to leverage complementary information, (5) Uncertainty quantification to know when predictions are unreliable. These advances will make automated segmentation safer and more reliable for clinical deployment.",
        slide_number=slide_num
    )
    
    # ==================== REFERENCES ====================
    slide_num += 1
    add_content_slide(
        "References (1/2)",
        [
            "1. Long, J., Shelhamer, E., & Darrell, T. (2015). Fully convolutional networks for semantic segmentation. CVPR.",
            "",
            "2. Ronneberger, O., Fischer, P., & Brox, T. (2015). U-net: Convolutional networks for biomedical image segmentation. MICCAI, 234-241.",
            "",
            "3. Milletari, F., Navab, N., & Ahmadi, S. A. (2016). V-net: Fully convolutional neural networks for volumetric medical image segmentation. 3DV, 565-571.",
            "",
            "4. Zhou, Z., Rahman Siddiquee, M. M., Tajbakhsh, N., & Liang, J. (2018). Unet++: A nested u-net architecture for medical image segmentation. MICCAI Workshops.",
            "",
            "5. Oktay, O., et al. (2018). Attention u-net: Learning where to look for the pancreas. MIDL.",
        ],
        "Part 1 of references covering foundational papers: FCN, U-Net, V-Net, U-Net++, and Attention U-Net. These papers established the core architectures used in medical image segmentation today.",
        slide_num
    )
    
    slide_num += 1
    add_content_slide(
        "References (2/2)",
        [
            "6. Isensee, F., Jaeger, P. F., Kohl, S. A., Petersen, J., & Maier-Hein, K. H. (2021). nnU-Net: a self-configuring method for deep learning-based biomedical image segmentation. Nature Methods, 18(2), 203-211.",
            "",
            "7. Chen, J., et al. (2021). TransUNet: Transformers make strong encoders for medical image segmentation. arXiv:2102.04306.",
            "",
            "8. Cao, H., et al. (2022). Swin-Unet: Unet-like pure transformer for medical image segmentation. ECCV Workshops.",
            "",
            "9. Zhou, H. Y., et al. (2023). nnFormer: Interleaved transformer for volumetric segmentation. arXiv:2109.03201.",
            "",
            "10. Medical Segmentation Decathlon: http://medicaldecathlon.com/",
        ],
        "Part 2 of references covering recent advances: nnU-Net (current state-of-the-art), transformer-based methods (TransUNet, Swin-UNET, nnFormer), and the Medical Segmentation Decathlon dataset used in our implementation.",
        slide_num
    )
    
    # Save presentation
    output_path = "/home/matt/im-seg/medical_image_segmentation_lecture.pptx"
    prs.save(output_path)
    print(f"\n✓ Presentation created successfully!")
    print(f"  File: {output_path}")
    print(f"  Total slides: {slide_num}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
