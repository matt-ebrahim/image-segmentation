#!/usr/bin/env python3
"""
Convert HTML presentation to PowerPoint while preserving style and formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def parse_html_presentation(html_path):
    """Parse HTML file and extract slides"""
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    slides = soup.find_all('div', class_='slide')
    
    return slides

def create_powerpoint_from_html(html_path, output_path):
    """Convert HTML presentation to PowerPoint"""
    
    # Parse HTML
    html_slides = parse_html_presentation(html_path)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio
    
    # Define color scheme matching HTML
    BLUE = RGBColor(0, 0, 255)
    DARK_TEXT = RGBColor(44, 62, 80)
    LIGHT_BLUE_BG = RGBColor(230, 230, 255)
    CREAM_BG = RGBColor(254, 249, 231)
    PINK_BG = RGBColor(253, 242, 248)
    GREEN_BG = RGBColor(232, 248, 245)
    PURPLE_BG = RGBColor(245, 238, 248)
    RED_BG = RGBColor(253, 234, 234)
    GRAY_BG = RGBColor(245, 246, 250)
    ORANGE_BG = RGBColor(254, 243, 226)
    
    # Border colors
    YELLOW_BORDER = RGBColor(244, 208, 63)
    PINK_BORDER = RGBColor(236, 72, 153)
    GREEN_BORDER = RGBColor(26, 188, 156)
    PURPLE_BORDER = RGBColor(155, 89, 182)
    RED_BORDER = RGBColor(231, 76, 60)
    GRAY_BORDER = RGBColor(127, 140, 141)
    ORANGE_BORDER = RGBColor(230, 126, 34)
    
    def add_blank_slide():
        """Add a blank slide"""
        return prs.slides.add_slide(prs.slide_layouts[6])
    
    def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
        """Add a text box to slide"""
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.text = text
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = color
            paragraph.font.bold = bold
            paragraph.alignment = alignment
            paragraph.font.name = 'Raleway'
        
        return textbox
    
    def add_colored_box(slide, left, top, width, height, bg_color, border_color, border_width=3):
        """Add a colored box shape"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        
        # Set fill
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        
        # Set border
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
        
        # Remove shadow
        shape.shadow.inherit = False
        
        return shape
    
    def extract_text_content(element):
        """Extract text content from HTML element, preserving structure"""
        if element is None:
            return ""
        
        text_parts = []
        for child in element.descendants:
            if child.name == 'br':
                text_parts.append('\n')
            elif child.string and child.string.strip():
                text_parts.append(child.string)
        
        return ' '.join(text_parts)
    
    def process_slide_content(slide_element):
        """Extract structured content from a slide"""
        content = {
            'title': '',
            'boxes': [],
            'text_blocks': [],
            'tables': [],
            'slide_number': ''
        }
        
        # Extract title
        title_elem = slide_element.find('h2', class_='slide-title')
        if title_elem:
            content['title'] = title_elem.get_text(strip=True)
        
        # Check for centered title slide (slide 1)
        h1_elem = slide_element.find('h1')
        if h1_elem:
            content['title'] = h1_elem.get_text(strip=True)
            content['is_title_slide'] = True
            # Get subtitle
            subtitle_elems = slide_element.find_all('p', style=lambda x: x and 'font-size: 20px' in x)
            if subtitle_elems:
                content['subtitle'] = subtitle_elems[0].get_text(strip=True)
        else:
            content['is_title_slide'] = False
        
        # Extract colored boxes
        box_classes = ['blue-box', 'cream-box', 'pink-box', 'green-box', 'purple-box', 'red-box', 'gray-box', 'orange-box']
        for box_class in box_classes:
            boxes = slide_element.find_all('div', class_=box_class)
            for box in boxes:
                box_content = {
                    'type': box_class,
                    'text': box.get_text(separator='\n', strip=True)
                }
                content['boxes'].append(box_content)
        
        # Extract tables
        tables = slide_element.find_all('table')
        for table in tables:
            table_data = []
            rows = table.find_all('tr')
            for row in rows:
                cells = row.find_all(['th', 'td'])
                row_data = [cell.get_text(strip=True) for cell in cells]
                table_data.append(row_data)
            content['tables'].append(table_data)
        
        # Extract slide number
        slide_num_elem = slide_element.find('span', class_='slide-number')
        if slide_num_elem:
            content['slide_number'] = slide_num_elem.get_text(strip=True)
        
        # Extract figure placeholders
        fig_placeholders = slide_element.find_all('div', class_='figure-placeholder')
        for fig in fig_placeholders:
            content['text_blocks'].append({
                'type': 'figure',
                'text': fig.get_text(separator='\n', strip=True)
            })
        
        # Extract gradient banners
        grad_banners = slide_element.find_all('div', class_='gradient-banner')
        for banner in grad_banners:
            content['text_blocks'].append({
                'type': 'gradient-banner',
                'text': banner.get_text(separator='\n', strip=True)
            })
        
        return content
    
    print(f"Converting {len(html_slides)} slides from HTML to PowerPoint...")
    
    # Process each slide
    for idx, html_slide in enumerate(html_slides, 1):
        slide = add_blank_slide()
        
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        content = process_slide_content(html_slide)
        
        # Handle title slide (Slide 1)
        if content.get('is_title_slide'):
            # Centered title slide
            title_box = add_text_box(slide, 1, 2, 8, 1, content['title'], 
                                    font_size=36, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
            
            if 'subtitle' in content:
                subtitle_box = add_text_box(slide, 1, 3, 8, 0.5, content['subtitle'],
                                           font_size=20, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
            
            # Add date if present
            date_text = "Oct 18 2025"
            date_box = add_text_box(slide, 1, 3.8, 8, 0.3, date_text,
                                   font_size=16, color=RGBColor(136, 136, 136), alignment=PP_ALIGN.CENTER)
        
        else:
            # Regular slide with title
            if content['title']:
                # Add title
                title_box = add_text_box(slide, 0.5, 0.3, 9, 0.6, content['title'],
                                        font_size=26, bold=True, color=BLUE)
                
                # Add underline
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.5), Inches(0.95), Inches(9), Inches(0.02)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = BLUE
                line.line.fill.background()
            
            # Add colored boxes and content
            current_y = 1.2
            
            # Process boxes
            for box in content['boxes']:
                box_type = box['type']
                text = box['text']
                
                # Determine colors based on box type
                if 'blue' in box_type:
                    bg_color, border_color = LIGHT_BLUE_BG, BLUE
                elif 'cream' in box_type:
                    bg_color, border_color = CREAM_BG, YELLOW_BORDER
                elif 'pink' in box_type:
                    bg_color, border_color = PINK_BG, PINK_BORDER
                elif 'green' in box_type:
                    bg_color, border_color = GREEN_BG, GREEN_BORDER
                elif 'purple' in box_type:
                    bg_color, border_color = PURPLE_BG, PURPLE_BORDER
                elif 'red' in box_type:
                    bg_color, border_color = RED_BG, RED_BORDER
                elif 'gray' in box_type:
                    bg_color, border_color = GRAY_BG, GRAY_BORDER
                elif 'orange' in box_type:
                    bg_color, border_color = ORANGE_BG, ORANGE_BORDER
                else:
                    bg_color, border_color = RGBColor(240, 240, 240), DARK_TEXT
                
                # Calculate box height based on text length
                lines = text.count('\n') + len(text) // 80
                box_height = min(0.15 * lines + 0.3, 1.5)
                
                # Add colored box shape
                box_shape = add_colored_box(slide, 0.7, current_y, 8.6, box_height, 
                                           bg_color, border_color)
                
                # Add text on top of box
                text_box = slide.shapes.add_textbox(Inches(0.8), Inches(current_y + 0.05),
                                                   Inches(8.4), Inches(box_height - 0.1))
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.text = text
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = DARK_TEXT
                    paragraph.font.name = 'Raleway'
                    paragraph.space_after = Pt(4)
                
                current_y += box_height + 0.15
                
                # Stop if we run out of vertical space
                if current_y > 4.8:
                    break
            
            # Add tables if present
            if content['tables'] and current_y < 4.0:
                for table_data in content['tables']:
                    if not table_data:
                        continue
                    
                    rows = len(table_data)
                    cols = len(table_data[0]) if table_data else 0
                    
                    if rows == 0 or cols == 0:
                        continue
                    
                    # Calculate table dimensions
                    table_width = 8.0
                    table_height = min(rows * 0.35, 2.0)
                    
                    # Add table
                    table_shape = slide.shapes.add_table(rows, cols, Inches(1.0), Inches(current_y),
                                                        Inches(table_width), Inches(table_height))
                    table = table_shape.table
                    
                    # Fill table
                    for i, row_data in enumerate(table_data):
                        for j, cell_text in enumerate(row_data):
                            cell = table.cell(i, j)
                            cell.text = cell_text
                            
                            # Style header row
                            if i == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = BLUE
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                                    paragraph.font.size = Pt(11)
                                    paragraph.font.bold = True
                            else:
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(10)
                                    paragraph.font.color.rgb = DARK_TEXT
                    
                    current_y += table_height + 0.2
            
            # Add gradient banners
            for text_block in content['text_blocks']:
                if text_block['type'] == 'gradient-banner' and current_y < 4.5:
                    # Add gradient effect (approximate with solid blue for simplicity)
                    banner_shape = add_colored_box(slide, 0.7, current_y, 8.6, 0.5,
                                                  RGBColor(80, 80, 200), RGBColor(80, 80, 200), 0)
                    
                    # Add text
                    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(current_y + 0.05),
                                                       Inches(8.4), Inches(0.4))
                    text_frame = text_box.text_frame
                    text_frame.text = text_block['text']
                    
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                        paragraph.font.name = 'Raleway'
                    
                    current_y += 0.65
                
                elif text_block['type'] == 'figure' and current_y < 4.5:
                    # Add figure placeholder
                    fig_box = add_colored_box(slide, 1.0, current_y, 8.0, 0.6,
                                             RGBColor(250, 250, 255), BLUE, 2)
                    
                    text_box = slide.shapes.add_textbox(Inches(1.1), Inches(current_y + 0.1),
                                                       Inches(7.8), Inches(0.4))
                    text_frame = text_box.text_frame
                    text_frame.text = "[Figure: " + text_block['text'][:100] + "...]"
                    
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.color.rgb = BLUE
                        paragraph.font.italic = True
                        paragraph.alignment = PP_ALIGN.CENTER
                    
                    current_y += 0.75
        
        # Add slide number
        if content['slide_number']:
            slide_num_box = add_text_box(slide, 9, 5.2, 0.5, 0.3, content['slide_number'],
                                        font_size=11, color=RGBColor(153, 153, 153), alignment=PP_ALIGN.RIGHT)
        
        print(f"  Converted slide {idx}: {content.get('title', 'Title Slide')[:50]}...")
    
    # Save presentation
    prs.save(output_path)
    print(f"\n✓ PowerPoint created successfully!")
    print(f"  File: {output_path}")
    print(f"  Total slides: {len(html_slides)}")
    
    return output_path

if __name__ == "__main__":
    html_file = "/home/matt/im-seg/session07_medical_image_segmentation_classification.html"
    output_file = "/home/matt/im-seg/session07_medical_image_segmentation_classification.pptx"
    
    create_powerpoint_from_html(html_file, output_file)
